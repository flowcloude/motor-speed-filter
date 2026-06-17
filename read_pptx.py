import sys
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation(r'c:\Users\Administrator\Desktop\motor\MC SDK5.4_2020-03.pptx')
for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(text[:300])
    print()
print(f"Total slides: {len(prs.slides)}")
